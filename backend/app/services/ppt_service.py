from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import base64
import subprocess
import os
import tempfile
from typing import Dict, Any, Optional

# PDF imports removed as requested

from app.schemas.presentation import Presentation as PresentationSchema

class PPTService:
    """
    Service for handling PowerPoint presentation generation
    """
    
    @staticmethod
    def create_presentation(presentation_data: Dict[str, Any]) -> BytesIO:
        """
        Create a PowerPoint presentation from structured presentation data
        
        Args:
            presentation_data: Dictionary representation of a Presentation
            
        Returns:
            BytesIO object containing the PPTX file
        """
        # Create a new presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9 aspect ratio
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Define consistent color theme
        PRIMARY_COLOR = RGBColor(0x44, 0x72, 0xC4)  # Blue theme color
        SECONDARY_COLOR = RGBColor(0x5B, 0x9B, 0xD5)  # Lighter blue
        ACCENT_COLOR = RGBColor(0x70, 0xAD, 0x47)  # Green accent
        TEXT_COLOR = RGBColor(0x00, 0x00, 0x00)  # Black text
        BACKGROUND_COLOR = RGBColor(0xFF, 0xFF, 0xFF)  # White text
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        # Style the title slide with gradient background
        background = title_slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = PRIMARY_COLOR
        fill.gradient_stops[0].position = 0
        fill.gradient_angle = 45
        fill.gradient_stops[1].color.rgb = SECONDARY_COLOR
        fill.gradient_stops[1].position = 1
        
        # Set title
        title = title_slide.shapes.title
        title.text = presentation_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
        
        # Add content slides
        for slide_data in presentation_data['slides']:
            # Use a content slide layout
            content_slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Style the slide with the same gradient background as the title slide
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = PRIMARY_COLOR
            fill.gradient_stops[0].position = 0
            fill.gradient_angle = 45
            fill.gradient_stops[1].color.rgb = SECONDARY_COLOR
            fill.gradient_stops[1].position = 1
            
            # Add title to the slide with white text
            title = slide.shapes.title
            title.text = slide_data['title']
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            
            # Get the content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Content placeholder
                    content_placeholder = shape
                    break
            
            # Add bullet points with improved formatting
            if content_placeholder:
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Clear existing text
                
                # Add each bullet point
                for i, bullet_data in enumerate(slide_data['bullets']):
                    if i == 0:
                        # First bullet uses the first paragraph already in the text frame
                        p = text_frame.paragraphs[0]
                    else:
                        # Subsequent bullets need new paragraphs
                        p = text_frame.add_paragraph()
                    
                    p.text = bullet_data['text']
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text for bullets
                    p.level = 0  # Top level bullet
                    
                    # Add visual interest with different bullet characters and colors
                    # Unfortunately, python-pptx doesn't provide direct access to bullet character styling
                    # This would require manual XML manipulation which is beyond the scope of this implementation
            
            # Add slide number at the bottom right
            left = prs.slide_width - Inches(1)
            top = prs.slide_height - Inches(0.5)
            width = Inches(0.5)
            height = Inches(0.3)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = str(presentation_data['slides'].index(slide_data) + 1)
            tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White text
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Save presentation to BytesIO object
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        return pptx_bytes

    @staticmethod
    def get_presentation_base64(presentation_data: Dict[str, Any]) -> str:
        """
        Generate a base64-encoded string of the PowerPoint presentation
        
        Args:
            presentation_data: Dictionary with presentation data
            
        Returns:
            Base64 encoded string of the PPTX file
        """
        pptx_bytes = PPTService.create_presentation(presentation_data)
        encoded = base64.b64encode(pptx_bytes.getvalue()).decode('utf-8')
        return encoded

    # PDF creation method removed as requested
